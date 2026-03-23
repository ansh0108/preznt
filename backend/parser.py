import fitz  # PyMuPDF
import docx
import pptx
import requests
import re
import os

# Known non-skill words that bleed in from LinkedIn PDF header/contact sections
SKILL_BLOCKLIST = {
    "area", "linkedin", "github", "email", "phone", "location", "contact",
    "connections", "followers", "profile", "page", "of", "view", "full",
}


def is_valid_skill(s: str) -> bool:
    """Return True only if the string looks like a real skill."""
    s = s.strip()
    if not s or len(s) < 2 or len(s) > 50:
        return False
    # Skip if it contains digits (dates, years)
    if re.search(r'\d', s):
        return False
    # Skip if it looks like a full sentence
    if len(s.split()) > 5:
        return False
    # Skip if any word is in blocklist
    words = s.lower().split()
    if any(w in SKILL_BLOCKLIST for w in words):
        return False
    # Skip if it ends with "Area" or "India" etc (locations)
    if re.search(r'\b(area|india|illinois|chicago|mumbai|remote|usa|uk)\b', s.lower()):
        return False
    # Skip if it looks like a name (two capitalized words with no tech keywords)
    tech_keywords = re.compile(
        r'sql|python|r\b|tableau|power\s?bi|excel|azure|aws|ml|ai|data|analysis|'
        r'analytics|visualization|modeling|statistics|bi|etl|dbt|snowflake|'
        r'postgres|mysql|java|javascript|typescript|react|node|fastapi|flask|'
        r'machine learning|deep learning|nlp|llm|api|cloud|git|docker|spark|'
        r'hadoop|kafka|intelligence|management|systems|engineering|science|'
        r'consulting|research|strategy|management|communication|leadership',
        re.IGNORECASE
    )
    if not tech_keywords.search(s):
        # If no tech keyword, only accept if it's clearly a single-word skill
        if len(s.split()) > 1:
            return False
    return True


def parse_linkedin_pdf(filepath: str) -> dict:
    doc = fitz.open(filepath)
    full_text = ""
    for page in doc:
        full_text += page.get_text()
    doc.close()
    structured = parse_linkedin_structured(full_text)
    return {
        "type": "linkedin",
        "raw_text": full_text,
        "source": "LinkedIn Profile",
        "structured": structured
    }


def parse_linkedin_structured(text: str) -> dict:
    # ── PRE-PROCESS: join split date lines ────────────────────────────────────
    raw_lines = [l.strip() for l in text.split('\n') if l.strip()]
    lines = []
    month_re = re.compile(
        r'(Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|'
        r'Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)',
        re.IGNORECASE
    )
    i = 0
    while i < len(raw_lines):
        line = raw_lines[i]
        if (i + 1 < len(raw_lines) and
                month_re.search(line) and
                re.match(r'^\d{4}', raw_lines[i + 1])):
            lines.append(line + ' ' + raw_lines[i + 1])
            i += 2
        else:
            lines.append(line)
            i += 1

    lines = [l for l in lines if not re.match(
        r'page\s+\d+\s+of\s+\d+', l.lower())]

    HEADERS = {
        "top skills":     ["top skills"],
        "summary":        ["summary", "about"],
        "experience":     ["experience"],
        "education":      ["education"],
        "certifications": ["licenses & certifications", "certifications"],
        "projects":       ["projects"],
        "languages":      ["languages"],
    }

    def find_section(name) -> int:
        for idx, line in enumerate(lines):
            low = line.lower().strip()
            if any(low == v for v in HEADERS.get(name, [name])):
                return idx
        return -1

    def get_section(name, max_lines=120):
        start = find_section(name)
        if start == -1:
            return []
        other_starts = []
        for key in HEADERS:
            idx = find_section(key)
            if idx > start:
                other_starts.append(idx)
        end = min(other_starts) if other_starts else len(lines)
        return lines[start + 1: min(end, start + max_lines)]

    date_line_re = re.compile(
        r'(Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|'
        r'Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)'
        r'\s+\d{4}\s*[-–]\s*'
        r'(Present|Current|'
        r'Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|'
        r'Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?'
        r'|\d{4})',
        re.IGNORECASE
    )
    duration_re = re.compile(r'\(?\d+\s+(yr|year|mo|month)', re.IGNORECASE)
    inline_date_re = re.compile(r'[·•]\s*\(?(.*?\d{4}.*?)\)?$')

    # ── TOP SKILLS (strict filtering) ────────────────────────────────────────
    skill_lines = get_section("top skills", max_lines=8)  # hard cap at 8 lines
    skills = []
    for line in skill_lines:
        if ',' in line:
            for s in line.split(','):
                if is_valid_skill(s):
                    skills.append(s.strip())
        else:
            if is_valid_skill(line):
                skills.append(line.strip())

    # ── EXPERIENCE (date-anchor approach) ────────────────────────────────────
    exp_lines = get_section("experience")
    experience = []
    date_indices = [i for i, l in enumerate(
        exp_lines) if date_line_re.search(l)]

    for di, date_idx in enumerate(date_indices):
        raw_date = exp_lines[date_idx]
        dates = re.sub(r'\s*\([\d\s\w]+\)\s*$', '', raw_date).strip()

        back = []
        j = date_idx - 1
        while j >= 0 and len(back) < 3:
            candidate = exp_lines[j]
            if duration_re.search(candidate) and len(candidate) < 30:
                j -= 1
                continue
            if date_line_re.search(candidate):
                break
            if len(candidate) > 2 and len(candidate) < 120:
                back.insert(0, candidate)
            j -= 1

        title = back[-1] if back else ""
        company = back[-2] if len(back) >= 2 else ""

        if not title or len(title) > 80 or (title[0].islower() and len(title) > 30):
            continue

        next_date_idx = date_indices[di + 1] if di + \
            1 < len(date_indices) else len(exp_lines)
        forward_lines = exp_lines[date_idx + 1: next_date_idx]

        location = ""
        description_parts = []
        for fl in forward_lines:
            if duration_re.search(fl) and len(fl) < 30:
                continue
            if not location and len(fl) < 60 and re.search(r',|area|india|il\b|ny\b|ca\b|remote', fl.lower()):
                location = fl
            elif len(fl) > 25 and not date_line_re.search(fl):
                description_parts.append(fl)

        experience.append({
            "title": title,
            "company": company,
            "dates": dates,
            "location": location,
            # raw text, LLM will summarize
            "description": " ".join(description_parts[:6])
        })

    # ── EDUCATION ────────────────────────────────────────────────────────────
    edu_lines = get_section("education")
    education = []
    i = 0
    while i < len(edu_lines) and len(education) < 6:
        line = edu_lines[i]
        if len(line) < 3 or date_line_re.search(line) or duration_re.search(line):
            i += 1
            continue
        school = line
        degree = ""
        dates = ""
        if i + 1 < len(edu_lines):
            next_line = edu_lines[i + 1]
            inline = inline_date_re.search(next_line)
            if inline:
                dates = inline.group(1).strip()
                degree = re.sub(r'\s*[·•]\s*\(.*', '', next_line).strip()
                i += 2
            elif date_line_re.search(next_line):
                dates = next_line
                i += 2
            else:
                degree = next_line
                if i + 2 < len(edu_lines):
                    date_candidate = edu_lines[i + 2]
                    inline2 = inline_date_re.search(date_candidate)
                    if inline2:
                        dates = inline2.group(1).strip()
                        degree = re.sub(r'\s*[·•]\s*\(.*', '', degree).strip()
                        i += 3
                    elif date_line_re.search(date_candidate) or re.search(r'\d{4}', date_candidate):
                        dates = re.sub(r'[()·]', '', date_candidate).strip()
                        i += 3
                    else:
                        i += 2
                else:
                    i += 2
        else:
            i += 1

        dates = re.sub(r'[()·]', '', dates).strip()
        degree = re.sub(r'\s*[·•].*$', '', degree).strip()

        if school and len(school) > 3:
            education.append(
                {"school": school, "degree": degree, "dates": dates})

    # ── SUMMARY ──────────────────────────────────────────────────────────────
    summary_lines = get_section("summary", max_lines=12)
    summary = " ".join(summary_lines[:8]).strip()

    print(f"\n[LinkedIn Parser] Lines: {len(lines)}")
    print(
        f"[LinkedIn Parser] Experience: {len(experience)} | Education: {len(education)} | Skills: {len(skills)}")
    for exp in experience:
        print(f"  EXP: {exp['title']} @ {exp['company']} ({exp['dates']})")
    for edu in education:
        print(f"  EDU: {edu['school']} — {edu['degree']} ({edu['dates']})")
    print(f"  SKILLS: {skills}")

    return {
        "experience": experience,
        "education": education,
        "skills": skills[:20],
        "summary": summary[:1000]
    }


def parse_github_repos(github_urls: list) -> list:
    repos = []
    GITHUB_TOKEN = os.getenv("GITHUB_TOKEN", "")
    headers = {
        "Accept": "application/vnd.github.v3+json",
        **({"Authorization": f"token {GITHUB_TOKEN}"} if GITHUB_TOKEN else {})
    }
    for url in github_urls:
        try:
            match = re.search(r"github\.com/([^/]+)/([^/\s]+)", url)
            if not match:
                continue
            owner, repo_name = match.group(1), match.group(2).rstrip("/")
            resp = requests.get(
                f"https://api.github.com/repos/{owner}/{repo_name}", headers=headers, timeout=10)
            if resp.status_code != 200:
                continue
            data = resp.json()

            readme_text = ""
            readme_resp = requests.get(
                f"https://api.github.com/repos/{owner}/{repo_name}/readme", headers=headers, timeout=10)
            if readme_resp.status_code == 200:
                import base64
                raw = readme_resp.json().get("content", "")
                readme_text = base64.b64decode(
                    raw).decode("utf-8", errors="ignore")
                readme_text = re.sub(r'!\[.*?\]\(.*?\)', '', readme_text)
                readme_text = re.sub(
                    r'\[([^\]]+)\]\([^\)]+\)', r'\1', readme_text)
                readme_text = re.sub(r'#{1,6}\s*', '', readme_text)
                readme_text = re.sub(r'\n{3,}', '\n\n', readme_text).strip()

            repos.append({
                "type": "github",
                "name": data.get("name", ""),
                "description": data.get("description") or "",
                "language": data.get("language") or "",
                "stars": data.get("stargazers_count", 0),
                "forks": data.get("forks_count", 0),
                "topics": data.get("topics", []),
                "url": url,
                "readme": readme_text[:4000],
                "source": f"GitHub: {repo_name}"
            })
        except Exception as e:
            print(f"[GitHub] Error: {url}: {e}")
    return repos


def parse_pdf_file(filepath: str, filename: str) -> dict:
    doc = fitz.open(filepath)
    full_text = "".join(page.get_text() for page in doc)
    doc.close()
    return {"type": "document", "raw_text": full_text[:8000], "source": filename}


def parse_docx_file(filepath: str, filename: str) -> dict:
    doc = docx.Document(filepath)
    return {"type": "document", "raw_text": "\n".join(p.text for p in doc.paragraphs)[:8000], "source": filename}


def parse_pptx_file(filepath: str, filename: str) -> dict:
    prs = pptx.Presentation(filepath)
    texts = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(
        shape, "text")]
    return {"type": "document", "raw_text": "\n".join(texts)[:8000], "source": filename}


def parse_txt_file(filepath: str, filename: str) -> dict:
    with open(filepath, "r", errors="ignore") as f:
        return {"type": "document", "raw_text": f.read()[:8000], "source": filename}


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def prepare_documents(data_items: list) -> list:
    documents = []
    for item in data_items:
        text = item.get("raw_text", "")
        if not text and item.get("type") == "github":
            text = f"""Project: {item.get('name', '')}
Description: {item.get('description', '')}
Language: {item.get('language', '')}
Topics: {', '.join(item.get('topics', []))}
README: {item.get('readme', '')}""".strip()
        for chunk in chunk_text(text):
            documents.append({"text": chunk, "source": item.get(
                "source", "Unknown"), "type": item.get("type", "document")})
    return documents
